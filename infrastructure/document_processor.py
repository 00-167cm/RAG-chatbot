"""
📄 ドキュメント処理
    各種ファイルからテキストを抽出し、チャンクに分割する
    
【役割】
- ファイルの読み込み（PDF / HTML / Excel / Word / PowerPoint）
- テキスト抽出（通常 + OCRフォールバック）
- チャンク分割（RAG用に適切なサイズに）
- メタデータの付与（ファイル名、ページ番号など）
"""
from typing import List, Dict, Any
from pathlib import Path

# PDF処理用ライブラリ
# PyMuPDF
import fitz 

# OCR用ライブラリ
from pdf2image import convert_from_path
import pytesseract

# HTML処理用ライブラリ
from bs4 import BeautifulSoup

# LangChainのテキスト分割
from langchain_text_splitters import RecursiveCharacterTextSplitter

from config.settings import CHUNK_SIZE, CHUNK_OVERLAP


class DocumentProcessor:
    """
    ドキュメント処理クラス
    各種ファイルをテキスト化してチャンクに分割する
    
    【このクラスが持つデータ】
    - self.chunk_size: 1チャンクの最大文字数
    - self.chunk_overlap: チャンク間の重複文字数
    - self.text_splitter: テキスト分割器（LangChain）
    """
    
    def __init__(
        self,
        chunk_size: int = CHUNK_SIZE,
        chunk_overlap: int = CHUNK_OVERLAP
    ):
        """
        ドキュメント処理の初期化
        
        Args:
            chunk_size: 1チャンクの最大文字数（デフォルト500文字）
            chunk_overlap: チャンク間の重複文字数（デフォルト100文字）
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            separators=["\n\n", "\n", "。", "、", " ", ""]
        )
    
    # ========================================
    # PDF処理
    # ========================================
    
    def extract_text_from_pdf(self, pdf_path: str) -> List[Dict[str, Any]]:
        """
        PDFからテキストを抽出（ページごと）
        通常抽出で0文字の場合、OCRにフォールバック
        
        Args:
            pdf_path: PDFファイルのパス
        
        Returns:
            ページごとのテキストとメタデータのリスト
        """
        pages = []
        file_name = Path(pdf_path).name
        
        try:
            # ========================================
            # Step 1: PyMuPDFで通常のテキスト抽出を試みる
            # ========================================
            doc = fitz.open(pdf_path)
            total_text_length = 0
            
            for page_num, page in enumerate(doc, start=1):
                text = page.get_text()
                total_text_length += len(text.strip())
                
                if text.strip():
                    pages.append({
                        "text": text,
                        "metadata": {
                            "source": file_name,
                            "page": page_num
                        }
                    })
            
            doc.close()
            
            # テキストが取れていればそのまま返す
            if total_text_length > 0:
                print(f"✅ PDF読み込み完了（テキスト抽出）: {file_name} ({len(pages)}ページ)")
                return pages
            
            # ========================================
            # Step 2: テキスト0文字 → OCRフォールバック
            # ========================================
            print(f"🔄 テキスト抽出0文字のためOCRに切り替え: {file_name}")
            pages = self._extract_text_with_ocr(pdf_path, file_name)
            
            if pages:
                print(f"✅ PDF読み込み完了（OCR）: {file_name} ({len(pages)}ページ)")
            else:
                print(f"⚠️ OCRでもテキストを抽出できませんでした: {file_name}")
            
            return pages
            
        except Exception as e:
            print(f"❌ PDF読み込みエラー: {pdf_path}")
            print(f"   エラー内容: {e}")
            return []
    
    def _extract_text_with_ocr(self, pdf_path: str, file_name: str) -> List[Dict[str, Any]]:
        """
        OCRでPDFからテキストを抽出（画像ベースPDF用）
        
        処理フロー:
            PDF → [poppler/pdf2image] → 画像 → [Tesseract OCR] → テキスト
        
        Args:
            pdf_path: PDFファイルのパス
            file_name: ファイル名（メタデータ用）
        
        Returns:
            ページごとのテキストとメタデータのリスト
        """
        pages = []
        
        try:
            # PDFを画像に変換（ページごと）
            # dpi=300 で高解像度変換（OCR精度向上のため）
            images = convert_from_path(pdf_path, dpi=300)
            
            print(f"   📸 {len(images)}ページを画像に変換しました")
            
            for page_num, image in enumerate(images, start=1):
                # Tesseract OCRで日本語テキスト抽出
                # lang="jpn" で日本語モード
                text = pytesseract.image_to_string(image, lang="jpn")
                
                if text.strip():
                    pages.append({
                        "text": text,
                        "metadata": {
                            "source": file_name,
                            "page": page_num
                        }
                    })
                
                # 進捗表示（10ページごと）
                if page_num % 10 == 0:
                    print(f"   📄 OCR処理中... {page_num}/{len(images)}ページ")
            
            return pages
            
        except Exception as e:
            print(f"❌ OCR処理エラー: {file_name}")
            print(f"   エラー内容: {e}")
            return []

    def process_pdf(self, pdf_path: str) -> List[Dict[str, Any]]:
        """PDFを処理してチャンクに分割"""
        pages = self.extract_text_from_pdf(pdf_path)
        if not pages:
            return []
        return self.split_into_chunks(pages)

    # ========================================
    # HTML処理
    # ========================================

    def extract_text_from_html(self, html_path: str) -> List[Dict[str, Any]]:
        """
        HTMLからテキストを抽出
        """
        file_name = Path(html_path).name

        try:
            with open(html_path, 'r', encoding='utf-8') as f:
                html_content = f.read()

            soup = BeautifulSoup(html_content, 'html.parser')

            for script in soup(['script', 'style']):
                script.decompose()

            text = soup.get_text()

            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)

            if text.strip():
                print(f"✅ HTML読み込み完了: {file_name}")
                return [{
                    "text": text,
                    "metadata": {
                        "source": file_name,
                        "page": 1
                    }
                }]
            else:
                print(f"⚠️ テキストが抽出できませんでした: {file_name}")
                return []

        except Exception as e:
            print(f"❌ HTML読み込みエラー: {html_path}")
            print(f"   エラー内容: {e}")
            return []

    def process_html(self, html_path: str) -> List[Dict[str, Any]]:
        """HTMLを処理してチャンクに分割"""
        pages = self.extract_text_from_html(html_path)
        if not pages:
            return []
        return self.split_into_chunks(pages)

    # ========================================
    # 🆕 Excel処理
    # ========================================

    def extract_text_from_excel(self, excel_path: str) -> List[Dict[str, Any]]:
        """
        Excelからテキストを抽出（シートごと）
        
        Args:
            excel_path: Excelファイルのパス
        
        Returns:
            シートごとのテキストとメタデータのリスト
        """
        file_name = Path(excel_path).name

        try:
            from openpyxl import load_workbook

            wb = load_workbook(excel_path, data_only=True)
            pages = []

            for sheet_num, sheet_name in enumerate(wb.sheetnames, start=1):
                ws = wb[sheet_name]
                rows = []

                for row in ws.iter_rows(values_only=True):
                    cell_values = [str(cell) if cell is not None else "" for cell in row]
                    row_text = " ".join(cell_values).strip()
                    if row_text:
                        rows.append(row_text)

                text = "\n".join(rows)

                if text.strip():
                    pages.append({
                        "text": text,
                        "metadata": {
                            "source": file_name,
                            "page": sheet_num
                        }
                    })

            wb.close()

            if pages:
                print(f"✅ Excel読み込み完了: {file_name} ({len(pages)}シート)")
            else:
                print(f"⚠️ テキストが抽出できませんでした: {file_name}")

            return pages

        except Exception as e:
            print(f"❌ Excel読み込みエラー: {excel_path}")
            print(f"   エラー内容: {e}")
            return []

    def process_excel(self, excel_path: str) -> List[Dict[str, Any]]:
        """Excelを処理してチャンクに分割"""
        pages = self.extract_text_from_excel(excel_path)
        if not pages:
            return []
        return self.split_into_chunks(pages)

    # ========================================
    # 🆕 Word処理
    # ========================================

    def extract_text_from_word(self, word_path: str) -> List[Dict[str, Any]]:
        """
        Wordからテキストを抽出
        
        Args:
            word_path: Wordファイルのパス
        
        Returns:
            テキストとメタデータのリスト
        """
        file_name = Path(word_path).name

        try:
            from docx import Document

            doc = Document(word_path)
            paragraphs = []

            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)

            # テーブル内のテキストも抽出
            for table in doc.tables:
                for row in table.rows:
                    cell_values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cell_values:
                        paragraphs.append(" ".join(cell_values))

            text = "\n".join(paragraphs)

            if text.strip():
                print(f"✅ Word読み込み完了: {file_name}")
                return [{
                    "text": text,
                    "metadata": {
                        "source": file_name,
                        "page": 1
                    }
                }]
            else:
                print(f"⚠️ テキストが抽出できませんでした: {file_name}")
                return []

        except Exception as e:
            print(f"❌ Word読み込みエラー: {word_path}")
            print(f"   エラー内容: {e}")
            return []

    def process_word(self, word_path: str) -> List[Dict[str, Any]]:
        """Wordを処理してチャンクに分割"""
        pages = self.extract_text_from_word(word_path)
        if not pages:
            return []
        return self.split_into_chunks(pages)

    # ========================================
    # 🆕 PowerPoint処理
    # ========================================

    def extract_text_from_pptx(self, pptx_path: str) -> List[Dict[str, Any]]:
        """
        PowerPointからテキストを抽出（スライドごと）
        
        Args:
            pptx_path: PowerPointファイルのパス
        
        Returns:
            スライドごとのテキストとメタデータのリスト
        """
        file_name = Path(pptx_path).name

        try:
            from pptx import Presentation

            prs = Presentation(pptx_path)
            pages = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                texts = []

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                texts.append(paragraph.text)

                    if shape.has_table:
                        for row in shape.table.rows:
                            cell_values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cell_values:
                                texts.append(" ".join(cell_values))

                text = "\n".join(texts)

                if text.strip():
                    pages.append({
                        "text": text,
                        "metadata": {
                            "source": file_name,
                            "page": slide_num
                        }
                    })

            if pages:
                print(f"✅ PowerPoint読み込み完了: {file_name} ({len(pages)}スライド)")
            else:
                print(f"⚠️ テキストが抽出できませんでした: {file_name}")

            return pages

        except Exception as e:
            print(f"❌ PowerPoint読み込みエラー: {pptx_path}")
            print(f"   エラー内容: {e}")
            return []

    def process_pptx(self, pptx_path: str) -> List[Dict[str, Any]]:
        """PowerPointを処理してチャンクに分割"""
        pages = self.extract_text_from_pptx(pptx_path)
        if not pages:
            return []
        return self.split_into_chunks(pages)

    # ========================================
    # 共通処理
    # ========================================

    def split_into_chunks(
        self,
        pages: List[Dict[str, Any]]
    ) -> List[Dict[str, Any]]:
        """
        ページテキストをチャンクに分割
        """
        chunks = []
        chunk_index = 0
        
        for page in pages:
            page_chunks = self.text_splitter.split_text(page["text"])
            
            for chunk_text in page_chunks:
                chunks.append({
                    "text": chunk_text,
                    "metadata": {
                        **page["metadata"],
                        "chunk_index": chunk_index
                    }
                })
                chunk_index += 1
        
        print(f"✅ チャンク分割完了: {len(chunks)}チャンク")
        return chunks

    def process_directory(self, directory_path: str) -> List[Dict[str, Any]]:
        """
        ディレクトリ内の全対応ファイルを処理
        """
        all_chunks = []
        directory = Path(directory_path)

        if not directory.exists():
            print(f"❌ ディレクトリが存在しません: {directory_path}")
            return []

        pdf_files = list(directory.glob("*.pdf"))
        html_files = list(directory.glob("*.html"))
        excel_files = list(directory.glob("*.xlsx"))
        word_files = list(directory.glob("*.docx"))
        pptx_files = list(directory.glob("*.pptx"))

        total_files = len(pdf_files) + len(html_files) + len(excel_files) + len(word_files) + len(pptx_files)

        if total_files == 0:
            print(f"⚠️ 対応ファイルが見つかりません: {directory_path}")
            return []

        print(f"📁 検出ファイル: PDF {len(pdf_files)}件, HTML {len(html_files)}件, "
              f"Excel {len(excel_files)}件, Word {len(word_files)}件, PowerPoint {len(pptx_files)}件")

        for pdf_file in pdf_files:
            chunks = self.process_pdf(str(pdf_file))
            all_chunks.extend(chunks)

        for html_file in html_files:
            chunks = self.process_html(str(html_file))
            all_chunks.extend(chunks)

        for excel_file in excel_files:
            chunks = self.process_excel(str(excel_file))
            all_chunks.extend(chunks)

        for word_file in word_files:
            chunks = self.process_word(str(word_file))
            all_chunks.extend(chunks)

        for pptx_file in pptx_files:
            chunks = self.process_pptx(str(pptx_file))
            all_chunks.extend(chunks)

        print(f"✅ 全ファイル処理完了: 合計{len(all_chunks)}チャンク")
        return all_chunks